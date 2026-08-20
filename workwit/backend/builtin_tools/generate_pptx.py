"""PPT 生成：专业级渲染引擎（对标 WorkBuddy pptx-generator 质量）。

质量标准：
  - 5 套专业配色主题（primary / secondary / accent / text / background 完整调色板）
  - 中文字体保障链（微软雅黑 -> SimHei -> PingFang SC -> Noto Sans CJK -> Arial）
  - 10 种幻灯片类型（cover / agenda / section / content / two_column / chart /
    table / quote / timeline / comparison / summary）
  - 渐变背景（封面/章节页）、阴影效果、圆角装饰、页码
  - 内容空值保护：无 content 时使用高质量默认模板
"""
import asyncio
import json
import os
import re
import tempfile
import time
import uuid

from builtin_tools._shared import _user_dir, _safe_id


# ============================================================
#  生成去重缓存（防止闭环反思重复生成同一主题的 PPT）
# ============================================================
# P2⑮ 缓存 key 鲁棒化：key = normalize(topic) + "||" + 当天日期。
#   - normalize(topic)：NFD 归一化 + 全部去空白 + 截断 50 字 + 小写。
#     解决"苏州一周天气概况" vs "苏州 1 周天气概况" vs "苏州一周天气概况（修订版）"
#     这种 LLM 反思时微调 topic 措辞导致 key 不一致、缓存失效的问题。
#   - 当天日期：避免跨天/重启后误命中旧版本（用户可能昨天生成的 PPT 与今天
#     同主题不同，跨天应允许重新生成）。
#   - TTL = 24 小时：同一天同一主题必复用，与旧版 5min 相比防止长任务中
#     缓存过早过期导致再次生成。
_recent_ppt_cache = {}  # cache_key -> (timestamp, filepath, had_content)
_CACHE_TTL = 86400  # 24 小时


def _normalize_topic_key(topic):
    """把 topic 归一化为缓存 key 的主题部分。"""
    import unicodedata as _ud
    t = (topic or "").strip()
    if not t:
        return ""
    # NFD 归一化：把组合字符拆开（中文不受影响，但防 emoji/全角字符混入）
    t = _ud.normalize("NFKD", t)
    # 全部去空白（中文/英文空格、tab、换行、零宽字符）
    t = re.sub(r"[\s\u00A0\u200B\u3000]+", "", t)
    # 截断防 key 过长
    t = t[:50]
    return t.lower()


def _ppt_cache_key(topic):
    """P2⑮ 鲁棒化缓存 key：归一化 topic + 当天日期。"""
    import datetime as _dt
    return _normalize_topic_key(topic) + "||" + _dt.date.today().isoformat()


# ============================================================
#  5 套专业配色主题（完整调色板，非单一主色）
# ============================================================
THEMES = {
    "business_blue": {
        "name": "商务蓝",
        "primary": (0x0B, 0x3D, 0x91),
        "secondary": (0x50, 0x8F, 0xC9),
        "accent": (0x00, 0xA3, 0xE0),
        "dark": (0x1A, 0x1A, 0x2E),
        "gray": (0x6B, 0x72, 0x80),
        "light_gray": (0xF0, 0xF4, 0xF8),
        "white": (0xFF, 0xFF, 0xFF),
        "success": (0x28, 0xA7, 0x45),
        "warning": (0xFF, 0xC1, 0x07),
        "danger": (0xDC, 0x35, 0x45),
        "cover_top": (0x0B, 0x3D, 0x91),
        "cover_bot": (0x1A, 0x53, 0x8E),
        "section_bg": (0x0B, 0x3D, 0x91),
    },
    "creative_purple": {
        "name": "创意紫",
        "primary": (0x6B, 0x21, 0xA8),
        "secondary": (0xA7, 0x5B, 0xD3),
        "accent": (0xD4, 0x9B, 0xFF),
        "dark": (0x1A, 0x1A, 0x2E),
        "gray": (0x6B, 0x72, 0x80),
        "light_gray": (0xF8, 0xF5, 0xFC),
        "white": (0xFF, 0xFF, 0xFF),
        "success": (0x20, 0xC9, 0x97),
        "warning": (0xFD, 0x7E, 0x14),
        "danger": (0xEF, 0x44, 0x44),
        "cover_top": (0x6B, 0x21, 0xA8),
        "cover_bot": (0x93, 0x3C, 0xD8),
        "section_bg": (0x6B, 0x21, 0xA8),
    },
    "tech_dark": {
        "name": "科技暗黑",
        "primary": (0x00, 0xD2, 0xFF),
        "secondary": (0x54, 0xA0, 0xFF),
        "accent": (0x00, 0xFF, 0x88),
        "dark": (0xE0, 0xE0, 0xE0),
        "gray": (0x94, 0xA3, 0xB8),
        "light_gray": (0x1A, 0x1A, 0x2E),
        "white": (0x0D, 0x11, 0x17),
        "success": (0x00, 0xFF, 0x88),
        "warning": (0xFF, 0xB8, 0x00),
        "danger": (0xFF, 0x55, 0x55),
        "cover_top": (0x0D, 0x11, 0x17),
        "cover_bot": (0x1A, 0x2A, 0x3A),
        "section_bg": (0x13, 0x18, 0x22),
    },
    "warm_orange": {
        "name": "活力橙",
        "primary": (0xE6, 0x5C, 0x00),
        "secondary": (0xF9, 0x8A, 0x10),
        "accent": (0xFF, 0xBA, 0x08),
        "dark": (0x2D, 0x2D, 0x2D),
        "gray": (0x6C, 0x75, 0x7D),
        "light_gray": (0xFF, 0xF8, 0xF0),
        "white": (0xFF, 0xFF, 0xFF),
        "success": (0x19, 0x8C, 0x58),
        "warning": (0xFD, 0x7E, 0x14),
        "danger": (0xDC, 0x35, 0x45),
        "cover_top": (0xE6, 0x5C, 0x00),
        "cover_bot": (0xF9, 0x8A, 0x10),
        "section_bg": (0xE6, 0x5C, 0x00),
    },
    "nature_green": {
        "name": "自然绿",
        "primary": (0x0F, 0x71, 0x41),
        "secondary": (0x2E, 0xCC, 0x71),
        "accent": (0x51, 0xCF, 0x66),
        "dark": (0x1A, 0x1A, 0x1A),
        "gray": (0x5F, 0x63, 0x68),
        "light_gray": (0xF0, 0xFB, 0xF4),
        "white": (0xFF, 0xFF, 0xFF),
        "success": (0x2E, 0xCC, 0x71),
        "warning": (0xFF, 0xCB, 0x05),
        "danger": (0xD9, 0x3A, 0x3A),
        "cover_top": (0x0F, 0x71, 0x41),
        "cover_bot": (0x1A, 0x8F, 0x5A),
        "section_bg": (0x0F, 0x71, 0x41),
    },
}

_DEFAULT_STYLE = "business_blue"
_FONT_STACK = [
    "微软雅黑", "Microsoft YaHei", "SimHei", "PingFang SC",
    "Noto Sans CJK SC", "WenQuanYi Micro Hei", "Arial Unicode MS", "Arial",
]


def _C(rgb_tuple):
    """Shortcut: tuple -> RGBColor."""
    from pptx.dml.color import RGBColor
    return RGBColor(*rgb_tuple)


def _set_font(run, size_pt=None, bold=None, color_rgb=None):
    """Set font on a run with Chinese font fallback chain."""
    fn = _FONT_STACK[0]
    run.font.name = fn
    # Set East Asian font for proper Chinese rendering in Office
    try:
        r = run._rElement
        rPr = r.get_or_add_rPr()
        from pptx.oxml.ns import qn
        from lxml import etree
        for tag in (qn("a:latin"), qn("a:ea")):
            for child in list(rPr):
                if child.tag.endswith("}latin") or child.tag.endswith("}ea") or child.tag.endswith("}cs"):
                    rPr.remove(child)
            rPr.insert(0, etree.SubElement(rPr, tag, typeface=fn))
    except Exception:
        pass
    if size_pt is not None:
        from pptx.util import Pt
        run.font.size = Pt(size_pt)
    if bold is not None:
        run.font.bold = bold
    if color_rgb is not None:
        run.font.color.rgb = color_rgb


def _add_page_number(slide, total, idx, theme):
    """Add page number at bottom-right."""
    from pptx.util import Inches, Pt
    tb = slide.shapes.add_textbox(Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.alignment = 2  # right
    r = p.add_run()
    r.text = "%d / %d" % (idx, total)
    _set_font(r, 10, color_rgb=_C(theme["gray"]))


def _add_rect(slide, left, top, w, h, fill_color):
    """Add a filled rectangle shape."""
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _C(fill_color) if isinstance(fill_color, tuple) else fill_color
    shape.line.fill.background()
    try:
        shape.adjustments[0] = 0.05
    except Exception:
        pass
    return shape


# ============================================================
#  Content parser
# ============================================================
def _parse_content(topic, content, style):
    """Parse user content into slides structure."""
    raw = (content or "").strip()
    if not raw:
        return None

    slides = {"style": style, "page_numbers": True,
              "slides": [{"type": "cover", "title": topic,
                          "subtitle": "AI 智能体自动生成"}]}

    # ---- Pre-scan: detect if content contains pipe-table syntax ----
    has_pipe_table = bool(re.search(r'^\|.+\|', raw, re.MULTILINE))

    if has_pipe_table:
        # Table mode: split into non-table blocks and table blocks
        # Tables are contiguous lines starting with |
        lines = raw.split("\n")
        current_block = []
        for ln in lines:
            stripped = ln.strip()
            if stripped.startswith("|"):
                current_block.append(stripped)
            else:
                # Non-table line: flush any accumulated table first
                if len(current_block) >= 2:
                    _append_table(slides, current_block)
                current_block = []
                # Process as regular text line
                if stripped:
                    _append_text_line(slides, stripped, topic)
        # Flush remaining table
        if len(current_block) >= 2:
            _append_table(slides, current_block)
    else:
        # Normal mode: split by --- or headings
        if "---" in raw:
            blocks = [b.strip() for b in raw.split("---") if b.strip()]
        else:
            blocks = [p.strip() for p in re.split(r"\n(?=#{1,3}\s)", raw) if p.strip()]

        for blk in blocks:
            lines_blk = [l.rstrip() for l in blk.split("\n") if l.strip()]
            if not lines_blk:
                continue
            # Check if this block looks like a table (all lines start with |)
            if all(l.startswith("|") for l in lines_blk) and len(lines_blk) >= 2:
                _append_table(slides, [l.strip() for l in lines_blk])
            else:
                title = re.sub(r"^#{1,3}\s*", "", lines_blk[0]).strip() or topic
                items = []
                for ln in lines_blk[1:]:
                    s = ln.strip()
                    if not s or s in ("---", "***", "___"):
                        continue
                    item = re.sub(r"^[-*+•]\s*", "", s).strip()
                    item = re.sub(r"^\d+[.)]\s*", "", item).strip()
                    if item and len(item) >= 2:
                        items.append(item)
                if not items and len(title) > 10:
                    parts = re.split(r"[;；。]\s*", title)
                    if len(parts) > 1:
                        title = parts[0]
                        items = [p.strip() for p in parts[1:] if p.strip() and len(p.strip()) >= 2]
                if not items:
                    items = [title]
                items = [it for it in items if len(it) >= 2 and it not in ("|", "\u2022", "-", "*")]
                if items:
                    slides["slides"].append({
                        "type": "content", "title": title[:45],
                        "items": items[:12]})

    return slides


def _append_table(slides, rows):
    """Parse pipe-delimited rows into a table slide."""
    table_rows = []
    for ln in rows:
        cells = [c.strip() for c in ln.split("|")]
        cells = [c for c in cells if c]
        if len(cells) >= 2 and not re.match(r"^[\s\-:]+$", "".join(cells)):
            table_rows.append(cells)
    if len(table_rows) >= 2:
        header = table_rows[0]
        data_rows = table_rows[1:]
        title = header[0] if header else "数据表格"
        slides["slides"].append({
            "type": "table", "title": title[:40],
            "header": header[:6], "rows": data_rows[:12]})


# ============================================================
#  噪声过滤 + 有效内容统计（防分页爆炸、稀烂 content 回退默认模板）
# ============================================================
_NOISE_CHARS = re.compile(r"^[\s\u2022·•—\-–*=_,.|:：；;!！?？\u2026/\\\u201c\u201d\u3002\u3001\u300a\u300b]+$")
# 噪声阈值：连续 N 行 ≤2 字符视为 LLM 反思/规划噪声
_NOISE_DENSE_RUN = 5


def _is_noise_line(text):
    """判断单行是否为噪声：空 / <2 字符 / 纯符号破折号点 / 整行 Markdown 装饰。"""
    t = (text or "").strip()
    if not t or len(t) < 2:
        return True
    if _NOISE_CHARS.match(t):
        return True
    return False


def _is_noise_dense(text):
    """检查文本是否含大量噪声行（连续 ≥5 行 ≤2 字符），若是则视为 LLM 反思/规划文本。
    用于 _auto_fill_content_from_history 兜底时剔除明显噪声资料。"""
    if not text:
        return True
    lines = text.split("\n")
    run = 0
    for ln in lines:
        if len(ln.strip()) <= 2:
            run += 1
            if run >= _NOISE_DENSE_RUN:
                return True
        else:
            run = 0
    return False


def _count_effective_content(slides):
    """统计有效内容页：cover/section/agenda/summary 不算；content/table/chart/two_column 算。
    用于判断"content 稀烂时是否回退默认模板"——有效内容页<2 时直接走默认模板 4 页+封面，
    避免「一行一页 117 页全是空破折号」式废稿。"""
    if not slides:
        return 0
    n = 0
    for s in slides.get("slides", []):
        t = s.get("type")
        if t in ("cover", "section", "agenda", "summary", "quote", "title"):
            continue
        if t == "content" and s.get("items"):
            n += 1
        elif t == "table" and (s.get("header") and s.get("rows")):
            n += 1
        elif t == "chart" and (s.get("data") or s.get("rows")):
            n += 1
        elif t == "two_column" and s.get("left") and s.get("right"):
            n += 1
        elif t == "timeline" and s.get("events"):
            n += 1
    return n


def _append_text_line(slides, text, topic):
    """把单行非空文本合并到上一页 items（短句像 bullet 时）或独立成页。

    设计要点（防分页爆炸）：
      - 噪声行（<2 字符 / 纯符号 / 破折号 / 单字符）直接跳过，不再独立成页
      - 短句（≤30 字符）若上一页是 content 且 items<12 → 合并到上一页 items（像 bullet）
      - 长句 / 标题 → 独立成页
    """
    text = text.strip()
    if _is_noise_line(text):
        return
    title = re.sub(r"^#{1,3}\s*", "", text).strip() or topic
    if _is_noise_line(title):
        return
    title = title[:45]
    # 短句合并到上一页 items（避免「一行一页 117 页」爆炸）
    if len(title) <= 30 and slides["slides"]:
        last = slides["slides"][-1]
        if (last.get("type") == "content" and
                isinstance(last.get("items"), list) and
                len(last["items"]) < 12):
            last["items"].append(title)
            return
    # 长句 / 标题：独立成页
    slides["slides"].append({"type": "content", "title": title, "items": [title]})


# ============================================================
#  阶段 2：结构化 slides JSON 入参（LLM 直接传结构数组，绕开 Markdown 启发式解析）
# ============================================================
_VALID_SLIDE_TYPES = {"cover", "section", "content", "agenda", "two_column",
                      "table", "chart", "summary", "quote", "timeline"}


def _build_slides_from_json(topic, slides_list, style):
    """把 LLM 直接传的结构化 slides JSON 转为 _build_and_save 期望的内部结构。

    入参示例（content 启发式解析猜不准时，用 slides 直接传最稳）：
        slides=[
            {"type":"cover","title":"苏州8月天气"},
            {"type":"content","title":"气温特征","items":["日间高温33-38℃","湿度70-90%","午后雷阵雨"]},
            {"type":"table","title":"周天气","header":["日期","最高温","最低温","降水"],
             "rows":[["8/1","35","26","雷阵雨"]]},
            {"type":"chart","title":"温度趋势","labels":["8/1","8/2"],"values":[35,36]},
            {"type":"summary","points":["建议携带雨具","注意防暑降温"]}
        ]

    校验/兜底：
      - 跳过非 dict / 未知 type 的项
      - 缺封面时自动补一个 cover 页
      - 必填字段缺失时安全降级（items=[] / rows=[] / data=[]）
    """
    out_slides = []
    has_cover = False
    for s in slides_list or []:
        if not isinstance(s, dict):
            continue
        t = s.get("type")
        if not t or t not in _VALID_SLIDE_TYPES:
            continue
        slide = {"type": t, "title": (s.get("title") or "").strip()[:45] or (topic or "PPT")}
        if t == "cover":
            slide["subtitle"] = (s.get("subtitle") or "AI 智能体自动生成").strip()[:60]
            has_cover = True
        elif t == "agenda":
            slide["topics"] = [str(x).strip() for x in (s.get("topics") or []) if str(x).strip()][:12]
        elif t == "content":
            slide["items"] = [str(x).strip() for x in (s.get("items") or []) if str(x).strip()][:12]
        elif t == "two_column":
            slide["left"] = [str(x).strip() for x in (s.get("left") or []) if str(x).strip()][:10]
            slide["right"] = [str(x).strip() for x in (s.get("right") or []) if str(x).strip()][:10]
        elif t == "table":
            slide["header"] = [str(x).strip() for x in (s.get("header") or []) if str(x).strip()][:6]
            slide["rows"] = [[str(c).strip() for c in row if str(c).strip()]
                             for row in (s.get("rows") or [])][:12]
        elif t == "chart":
            data = s.get("data") or []
            if not data and s.get("labels") and s.get("values"):
                data = [[l, v] for l, v in zip(s["labels"], s["values"])]
            slide["data"] = [[str(x[0]), float(x[1])] for x in data
                             if isinstance(x, (list, tuple)) and len(x) >= 2][:12]
        elif t == "summary":
            slide["points"] = [str(x).strip() for x in (s.get("points") or []) if str(x).strip()][:8]
        elif t == "quote":
            slide["quote"] = (s.get("quote") or s.get("text") or "").strip()[:200]
            slide["author"] = (s.get("author") or "").strip()[:40]
        elif t == "timeline":
            slide["events"] = [{"date": str(e.get("date") or "")[:20],
                                "title": str(e.get("title") or "")[:40]}
                               for e in (s.get("events") or []) if isinstance(e, dict)][:10]
        out_slides.append(slide)
    if not has_cover and topic:
        out_slides.insert(0, {"type": "cover", "title": topic, "subtitle": "AI 智能体自动生成"})
    if not out_slides:
        return None
    return {"style": style, "page_numbers": True, "slides": out_slides}


# ============================================================
#  Default template
# ============================================================
def _default_slides(topic, style):
    """High-quality default template when no content provided."""
    return {
        "style": style, "page_numbers": True,
        "slides": [
            {"type": "cover", "title": topic, "subtitle": "AI 智能体自动生成"},
            {"type": "agenda", "title": "目录",
             "topics": ["项目背景", "核心内容", "数据分析", "总结与展望"]},
            {"type": "section", "title": "01 项目背景"},
            {"type": "content", "title": "背景概述",
             "items": ["阐述本项目的背景与起源",
                      "说明项目的核心目标与预期价值",
                      "介绍相关的行业现状与痛点分析"]},
            {"type": "section", "title": "02 核心内容"},
            {"type": "content", "title": "主要内容",
             "items": ["围绕主题展开的第一项关键要点",
                      "围绕主题展开的第二项关键要点",
                      "围绕主题展开的第三项关键要点",
                      "实施路径与关键里程碑"]},
            {"type": "two_column", "title": "对比分析",
             "left_title": "方案 A", "left_items": ["优势 1", "优势 2"],
             "right_title": "方案 B", "right_items": ["优势 1", "优势 2"]},
            {"type": "section", "title": "03 数据分析"},
            {"type": "chart", "title": "关键指标趋势",
             "chart_type": "bar",
             "categories": ["Q1", "Q2", "Q3", "Q4"],
             "series": [{"name": "完成率 (%)", "values": [78, 85, 92, 96]}]},
            {"type": "summary", "title": "总结与展望",
             "points": ["本项目通过系统化方法实现了预期目标",
                       "关键数据指标均达到或超过既定标准",
                       "后续将持续优化并扩展应用场景"],
             "conclusion": "感谢观看"},
        ]
    }


# ============================================================
#  META
# ============================================================
META = {
    "name": "generate_ppt", "display_name": "PPT 生成", "category": "generation",
    "description": (
        "根据主题生成一份专业级 PPT 文件（python-pptx 后端渲染）。\n"
        "支持 5 套配色、10 种版式（封面/目录/章节/内容/双栏/表格/图表/摘要/引用/时间线）、"
        "真表格真图表（add_table / add_chart）、20 页硬上限。\n\n"

        "═══════════════════════════════════════════════════════════════\n"
        "【硬性规则·必须遵守】调用前请先完成「数据准备」\n"
        "═══════════════════════════════════════════════════════════════\n"
        "(1) 本工具是「无状态渲染器」——只看你传入的参数，不知道上下文里有什么数据。"
        "只传 topic 不传 content/slides → 只会得到空泛的「项目背景/核心内容/数据分析」通用模板，"
        "无用户真实数据。\n\n"
        "(2) **数据要传完整**——\n"
        "    · 表格必须用 Markdown 表格语法 `|列1|列2|` 填入完整数据行（不是一行表头 + 空 body）；\n"
        "    · 图表必须给列名 + 至少 3 个数据点；\n"
        "    · 时间线必须给 date + title 列表；\n"
        "    · ❌ 严禁只写「——」「·」「# 标题」这种占位文本（会生成 117 页全是空破折号的废稿）。\n\n"
        "(3) **分页用 `---` 或 `# 标题`**——把多主题拆成多页；同一页内用 `-` 列要点。\n\n"
        "(4) **同一个任务只调一次**——系统已对 topic+24h 做去重，重复调用不会产生新文件，"
        "只会浪费 step 配额。\n\n"

        "═══════════════════════════════════════════════════════════════\n"
        "【两种入参模式】选一种即可\n"
        "═══════════════════════════════════════════════════════════════\n"
        "**方式 A（推荐·更稳）：直接传 slides 结构数组**\n"
        "  slides=[{\"type\":\"table\",\"title\":\"...\",\"header\":[...],\"rows\":[[...]]}]\n"
        "  完全绕开 Markdown 启发式解析，LLM 不再「猜」，质量直接对齐 pptx-generator 范式。\n\n"
        "**方式 B（兼容·Markdown）：传 content**\n"
        "  content=\"## 气温特征\\n- 33-38℃\\n- 湿度70-90%\\n\\n## 周天气\\n| 日期 | 最高 |\\n|---|---|\\n| 8/1 | 35 |\"\n"
        "  用 `---` 或 `# 标题` 分页，表格用 `|列1|列2|` 语法。\n\n"

        "═══════════════════════════════════════════════════════════════\n"
        "【好 vs 坏 content 对比】\n"
        "═══════════════════════════════════════════════════════════════\n"
        "❌ 烂示例（会得到 100+ 页全是空破折号 / 单字符的废稿）：\n"
        "  content=\"——\\n——\\n·\\n·\\n**核心要点**\\n# 标题\\n—\\n• 结论\\n——\\n\"\n\n"
        "✅ 好示例 A（Markdown 方式）：\n"
        "  content=\"## 气温特征\\n- 日间高温33-38℃\\n- 湿度70-90%\\n\\n## 周天气\\n"
        "| 日期 | 最高 | 降水 |\\n|---|---|---|\\n| 8/1 | 35 | 雷阵雨 |\\n| 8/2 | 36 | 晴 |\"\n\n"
        "✅ 好示例 B（slides 结构数组，更稳）：\n"
        "  slides=[{\"type\":\"content\",\"title\":\"气温特征\",\"items\":[\"日间高温33-38℃\",\"湿度70-90%\"]},\n"
        "          {\"type\":\"table\",\"title\":\"周天气\",\"header\":[\"日期\",\"最高\",\"降水\"],"
        "\"rows\":[[\"8/1\",\"35\",\"雷阵雨\"]]},\n"
        "          {\"type\":\"summary\",\"points\":[\"建议携带雨具\",\"注意防暑降温\"]}]\n\n"

        "═══════════════════════════════════════════════════════════════\n"
        "【自动兜底】\n"
        "═══════════════════════════════════════════════════════════════\n"
        "· 没传 content/slides → 自动从聊天历史抽最近一条助手消息（≥80 字符、非噪声密集）；\n"
        "· content 解析后有效内容页 < 2 → 自动回退默认模板（4 页+封面）；\n"
        "· 页数 > 20 → 截断 + 告警（建议拆分为多份、多套配色，或改为 Word 报告）。\n\n"

        "当用户要求「做 PPT / 生成演示文稿 / 汇报材料 / 做成 PPT / 整理成 PPT」时调用。"
    ),
    "params": {"type": "object",
               "properties": {
                   "topic": {"type": "string",
                             "description": "PPT 主题/标题（必填）。封面会显示此文本。"},
                   "content": {"type": "string",
                              "description": (
                                  "正文内容（Markdown 方式）。支持 `---` 分页、`# 标题` 分页、"
                                  "Markdown 表格语法 `|列1|列2|`。表格必须填入完整数据行，"
                                  "图表必须给列名+数据点。填什么 PPT 就展示什么；"
                                  "不传则从聊天历史兜底或回退默认模板。")},
                   "slides": {"type": "array",
                              "description": (
                                  "【推荐】结构化 slides 数组，绕开 Markdown 启发式解析。"
                                  "每项形如 {\"type\":\"cover|section|content|agenda|two_column|"
                                  "table|chart|summary|quote|timeline\", "
                                  "\"title\":\"...\", "
                                  "\"items\":[...]|\"header\":[...],\"rows\":[[...]]|"
                                  "\"labels\":[...],\"values\":[...]|\"points\":[...]}。"
                                  "传了 slides 则优先于 content；content 与 slides 至少传一个。"),
                              "items": {"type": "object"}},
                   "style": {"type": "string",
                             "description": "配色主题：%s。默认 %s" % (
                                 ", ".join(THEMES.keys()), _DEFAULT_STYLE)},
               },
               "required": ["topic"]},
    "backend_type": "builtin", "handler": "generate_ppt",
    "trigger_words": "PPT,演示,幻灯片,汇报,宣讲,演示文稿,PowerPoint,pptx,整理成PPT,做成PPT,生成演示文稿",
}


# ============================================================
#  Rendering engine
# ============================================================
def _render_pptx(sd):
    """Render slides dict to python-pptx Presentation."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    style_name = sd.get("style", _DEFAULT_STYLE)
    theme = THEMES.get(style_name, THEMES[_DEFAULT_STYLE])
    slide_list = sd.get("slides", [])
    total = len(slide_list)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def ns():
        return prs.slides.add_slide(blank)

    def sf(run, sz=None, bold=None, clr=None):
        _set_font(run, sz, bold, clr)

    def rect(slide, l, t, w, h, fc):
        return _add_rect(slide, l, t, w, h, fc)

    # ---- cover ----
    def render_cover(s, idx):
        slide = ns()
        bg = slide.background.fill
        bg.gradient()
        bg.gradient_angle = 90
        bg.gradient_stops[0].color.rgb = _C(theme["cover_top"])
        bg.gradient_stops[1].color.rgb = _C(theme["cover_bot"])

        rect(slide, Inches(0.5), Inches(2.0), Inches(0.08), Inches(3.5), theme["accent"])

        tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.0), Inches(1.8))
        tf = tb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run(); r.text = s.get("title", "")
        sf(r, 48, True, _C(theme["white"]))

        sub = s.get("subtitle", "")
        if sub:
            stb = slide.shapes.add_textbox(Inches(1.1), Inches(4.5), Inches(10.0), Inches(0.8))
            sr = stb.text_frame.paragraphs[0].add_run(); sr.text = sub
            sf(sr, 20, False, _C((200, 210, 220)))

        rect(slide, Inches(1.0), Inches(5.5), Inches(4.0), Pt(4), theme["accent"])
        _add_page_number(slide, total, idx, theme)

    # ---- agenda ----
    def render_agenda(s, idx):
        slide = ns()
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.0), Inches(1.0))
        r = tb.text_frame.paragraphs[0].add_run(); r.text = s.get("title", "目录")
        sf(r, 36, True, _C(theme["primary"]))
        rect(slide, Inches(0.8), Inches(1.4), Inches(2.5), Pt(4), theme["primary"])

        topics = s.get("topics", [])
        for i, topic in enumerate(topics, 1):
            y = Inches(2.0 + (i - 1) * 0.95)
            circle = rect(slide, Inches(0.9), y, Inches(0.5), Inches(0.5), theme["primary"])
            cr = circle.text_frame.paragraphs[0]; cr.alignment = PP_ALIGN.CENTER
            crcr = cr.add_run(); crcr.text = str(i); sf(crcr, 18, True, _C(theme["white"]))

            tbox = slide.shapes.add_textbox(Inches(1.6), y + Inches(0.08), Inches(9.0), Inches(0.5))
            tr = tbox.text_frame.paragraphs[0].add_run(); tr.text = topic
            sf(tr, 22, False, _C(theme["dark"]))
            rect(slide, Inches(1.6), y + Inches(0.55), Inches(8.0), Pt(0.75), theme["light_gray"])

        _add_page_number(slide, total, idx, theme)

    # ---- section ----
    def render_section(s, idx):
        slide = ns()
        bg = slide.background.fill; bg.solid()
        bg.fore_color.rgb = _C(theme["section_bg"])
        rect(slide, Inches(0), Inches(0), Inches(0.25), Inches(7.5), theme["accent"])

        tb = slide.shapes.add_textbox(Inches(1.2), Inches(2.8), Inches(10.5), Inches(1.8))
        tf = tb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run(); r.text = s.get("title", "")
        sf(r, 44, True, _C(theme["white"]))
        rect(slide, Inches(1.2), Inches(4.8), Inches(4.0), Pt(3), theme["accent"])
        _add_page_number(slide, total, idx, theme)

    # ---- content ----
    def render_content(s, idx):
        slide = ns()
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12.0), Inches(1.0))
        r = tb.text_frame.paragraphs[0].add_run(); r.text = s.get("title", "")
        sf(r, 30, True, _C(theme["primary"]))
        rect(slide, Inches(0.7), Inches(1.35), Inches(0.08), Inches(0.5), theme["accent"])
        rect(slide, Inches(0.85), Inches(1.37), Inches(3.5), Pt(3.5), theme["secondary"])

        body_tb = slide.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(11.5), Inches(5.0))
        btf = body_tb.text_frame; btf.word_wrap = True
        for i, item in enumerate(s.get("items", [])):
            p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            p.space_before = Pt(8); p.space_after = Pt(6)
            dot_r = p.add_run(); dot_r.text = "\u25cf  "
            sf(dot_r, 14, False, _C(theme["accent"] if i < 3 else theme["secondary"]))
            tr = p.add_run(); tr.text = item
            sf(tr, 19, False, _C(theme["dark"]))
        _add_page_number(slide, total, idx, theme)

    # ---- two_column ----
    def render_two_column(s, idx):
        slide = ns()
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12.0), Inches(1.0))
        r = tb.text_frame.paragraphs[0].add_run(); r.text = s.get("title", "对比分析")
        sf(r, 30, True, _C(theme["primary"]))
        rect(slide, Inches(0.7), Inches(1.35), Inches(0.08), Inches(0.5), theme["accent"])
        rect(slide, Inches(0.85), Inches(1.37), Inches(3.5), Pt(3.5), theme["secondary"])

        rect(slide, Inches(0.7), Inches(1.9), Inches(5.8), Inches(5.0), theme["light_gray"])
        rect(slide, Inches(6.8), Inches(1.9), Inches(5.8), Inches(5.0), theme["light_gray"])

        for prefix, x_off, title_key, items_key in [("left", 1.0, "left_title", "left_items"),
                                                    ("right", 7.1, "right_title", "right_items")]:
            lt = slide.shapes.add_textbox(Inches(x_off), Inches(2.1), Inches(5.2), Inches(0.6))
            lr = lt.text_frame.paragraphs[0].add_run(); lr.text = s.get(title_key, "方案")
            sf(lr, 22, True, _C(theme["primary"]))
            lbody = slide.shapes.add_textbox(Inches(x_off), Inches(2.75), Inches(5.2), Inches(4.0))
            lbf = lbody.text_frame; lbf.word_wrap = True
            for j, item in enumerate(s.get(items_key, []) or []):
                lp = lbf.paragraphs[0] if j == 0 else lbf.add_paragraph()
                dr = lp.add_run(); dr.text = "\u2713  "; sf(dr, 14, False, _C(theme["success"]))
                ir = lp.add_run(); ir.text = item; sf(ir, 17, False, _C(theme["dark"]))
        _add_page_number(slide, total, idx, theme)

    # ---- chart ----
    def render_chart(s, idx):
        slide = ns()
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12.0), Inches(1.0))
        r = tb.text_frame.paragraphs[0].add_run(); r.text = s.get("title", "数据图表")
        sf(r, 30, True, _C(theme["primary"]))
        rect(slide, Inches(0.7), Inches(1.35), Inches(0.08), Inches(0.5), theme["accent"])
        rect(slide, Inches(0.85), Inches(1.37), Inches(3.5), Pt(3.5), theme["secondary"])

        cats = s.get("categories", [])
        series_list = s.get("series", [])
        if cats and series_list:
            ct_map = {"bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
                      "line": XL_CHART_TYPE.LINE, "pie": XL_CHART_TYPE.PIE}
            ct = ct_map.get(s.get("chart_type", "bar"), XL_CHART_TYPE.COLUMN_CLUSTERED)
            cd = CategoryChartData()
            cd.categories = cats
            for ser in series_list:
                cd.add_series(ser.get("name", ""), ser.get("values", []))
            gf = slide.shapes.add_chart(ct, Inches(0.8), Inches(1.8),
                                          Inches(11.7), Inches(5.2), cd)
            chart = gf.chart
            chart.has_legend = len(series_list) > 1
            if chart.has_legend:
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        else:
            hint = slide.shapes.add_textbox(Inches(3.0), Inches(3.5), Inches(7.0), Inches(1.0))
            hr = hint.text_frame.paragraphs[0]; hr.alignment = PP_ALIGN.CENTER
            hrr = hr.add_run(); hrr.text = "[ 图表数据待补充 ]"; sf(hrr, 18, None, _C(theme["gray"]))
        _add_page_number(slide, total, idx, theme)

    # ---- table ----
    def render_table(s, idx):
        slide = ns()
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12.0), Inches(1.0))
        r = tb.text_frame.paragraphs[0].add_run(); r.text = s.get("title", "数据表格")
        sf(r, 30, True, _C(theme["primary"]))
        rect(slide, Inches(0.7), Inches(1.35), Inches(0.08), Inches(0.5), theme["accent"])
        rect(slide, Inches(0.85), Inches(1.37), Inches(3.5), Pt(3.5), theme["secondary"])

        header = s.get("header", [])
        rows = s.get("rows", [])
        if header and rows:
            cols = min(len(header), 6)
            rcount = min(len(rows) + 1, 13)
            tbl = slide.shapes.add_table(rcount, cols,
                                          Inches(0.8), Inches(1.9),
                                          Inches(11.7), Inches(5.2)).table
            cw = int(11.7 / cols * 914400)
            for ci in range(cols):
                tbl.columns[ci].width = cw
            for ci, h in enumerate(header[:cols]):
                cell = tbl.cell(0, ci); cell.text = str(h)
                cell.fill.solid(); cell.fill.fore_color.rgb = _C(theme["primary"])
                cp = cell.text_frame.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
                cr = cp.runs[0] if cp.runs else cp.add_run(); cr.text = str(h)
                sf(cr, 15, True, _C(theme["white"]))
            for ri, row in enumerate(rows[:rcount - 1]):
                for ci, val in enumerate(row[:cols]):
                    cell = tbl.cell(ri + 1, ci); cell.text = str(val) if val else ""
                    if ri % 2 == 0:
                        cell.fill.solid(); cell.fill.fore_color.rgb = _C(theme["light_gray"])
                    cp = cell.text_frame.paragraphs[0]
                    cr = cp.runs[0] if cp.runs else cp.add_run()
                    cr.text = str(val) if val else ""; sf(cr, 14, False, _C(theme["dark"]))
        else:
            hint = slide.shapes.add_textbox(Inches(3.0), Inches(3.5), Inches(7.0), Inches(1.0))
            hr = hint.text_frame.paragraphs[0]; hr.alignment = PP_ALIGN.CENTER
            hrr = hr.add_run(); hrr.text = "[ 表格数据待补充 ]"; sf(hrr, 18, None, _C(theme["gray"]))
        _add_page_number(slide, total, idx, theme)

    # ---- quote ----
    def render_quote(s, idx):
        slide = ns()
        bg = slide.background.fill; bg.solid()
        bg.fore_color.rgb = _C(theme["light_gray"])

        qt = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(1.5), Inches(2.0))
        qr = qt.text_frame.paragraphs[0].add_run()
        qr.text = "\u201c"
        sf(qr, 120, True, _C((200, 210, 220, 128)))

        qb = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.0), Inches(3.0))
        qbf = qb.text_frame; qbf.word_wrap = True
        qbr = qbf.paragraphs[0].add_run()
        qbr.text = s.get("text", s.get("title", ""))
        sf(qbr, 28, False, _C(theme["dark"]))

        author = s.get("author", "")
        if author:
            atb = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(10.0), Inches(0.6))
            ap = atb.text_frame.paragraphs[0]; ap.alignment = PP_ALIGN.RIGHT
            ar = ap.add_run(); ar.text = "\u2014 " + author
            sf(ar, 16, None, _C(theme["gray"]))
        _add_page_number(slide, total, idx, theme)

    # ---- timeline ----
    def render_timeline(s, idx):
        slide = ns()
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12.0), Inches(1.0))
        r = tb.text_frame.paragraphs[0].add_run(); r.text = s.get("title", "发展历程")
        sf(r, 30, True, _C(theme["primary"]))
        rect(slide, Inches(0.7), Inches(1.35), Inches(0.08), Inches(0.5), theme["accent"])
        rect(slide, Inches(0.85), Inches(1.37), Inches(3.5), Pt(3.5), theme["secondary"])

        events = s.get("events", [])
        if events:
            rect(slide, Inches(1.0), Inches(4.0), Inches(11.3), Pt(3), theme["secondary"])
            n = len(events)
            sx = 11.3 / max(n, 1)
            for i, ev in enumerate(events):
                cx = Inches(1.0 + i * sx)
                node = rect(slide, cx - Inches(0.15), Inches(3.88),
                             Inches(0.3), Inches(0.3), theme["primary"])
                time_label = ev.get("time", "") if isinstance(ev, dict) else ""
                if time_label:
                    ttb = slide.shapes.add_textbox(cx - Inches(0.8), Inches(3.2),
                                                    Inches(1.6), Inches(0.5))
                    tp = ttb.text_frame.paragraphs[0]; tp.alignment = PP_ALIGN.CENTER
                    tr = tp.add_run(); tr.text = str(time_label); sf(tr, 13, True, _C(theme["primary"]))
                desc = ev.get("desc", "") if isinstance(ev, dict) else str(ev)
                dy = Inches(4.5) if i % 2 == 0 else Inches(2.3)
                dtb = slide.shapes.add_textbox(cx - Inches(0.8), dy, Inches(1.6), Inches(1.3))
                dtp = dtb.text_frame; dtp.word_wrap = True
                dtp.paragraphs[0].alignment = PP_ALIGN.CENTER
                dr = dtp.paragraphs[0].add_run(); dr.text = desc[:40]
                sf(dr, 12, False, _C(theme["dark"]))
        else:
            hint = slide.shapes.add_textbox(Inches(3.0), Inches(3.5), Inches(7.0), Inches(1.0))
            hr = hint.text_frame.paragraphs[0]; hr.alignment = PP_ALIGN.CENTER
            hrr = hr.add_run(); hrr.text = "[ 时间线事件待补充 ]"; sf(hrr, 18, None, _C(theme["gray"]))
        _add_page_number(slide, total, idx, theme)

    # ---- comparison ----
    def render_comparison(s, idx):
        slide = ns()
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12.0), Inches(1.0))
        r = tb.text_frame.paragraphs[0].add_run(); r.text = s.get("title", "多维对比")
        sf(r, 30, True, _C(theme["primary"]))
        rect(slide, Inches(0.7), Inches(1.35), Inches(0.08), Inches(0.5), theme["accent"])
        rect(slide, Inches(0.85), Inches(1.37), Inches(3.5), Pt(3.5), theme["secondary"])

        dims = s.get("dimensions", [])
        opts = s.get("options", [])
        scores = s.get("scores", [[]])
        if dims and opts:
            tbl = slide.shapes.add_table(len(dims) + 1, len(opts),
                                          Inches(1.5), Inches(1.9),
                                          Inches(10.3), Inches(5.2)).table
            for ci, opt in enumerate(opts):
                cell = tbl.cell(0, ci); cell.text = str(opt)
                cell.fill.solid(); cell.fill.fore_color.rgb = _C(theme["primary"])
                cp = cell.text_frame.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
                cr = cp.runs[0] if cp.runs else cp.add_run(); cr.text = str(opt)
                sf(cr, 15, True, _C(theme["white"]))
            for ri, dim in enumerate(dims):
                dcell = tbl.cell(ri + 1, 0); dcell.text = str(dim)
                dcell.fill.solid(); dcell.fill.fore_color.rgb = _C(theme["light_gray"])
                dp = dcell.text_frame.paragraphs[0]
                dr = dp.runs[0] if dp.runs else dp.add_run(); dr.text = str(dim)
                sf(dr, 14, True, _C(theme["dark"]))
                for ci in range(len(opts)):
                    cell = tbl.cell(ri + 1, ci)
                    val = ""
                    if ri < len(scores) and ci < len(scores[ri]):
                        val = str(scores[ri][ci])
                    cell.text = val
                    if ri % 2 == 0:
                        cell.fill.solid(); cell.fill.fore_color.rgb = _C(theme["white"])
                    cp = cell.text_frame.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
                    cr = cp.runs[0] if cp.runs else cp.add_run(); cr.text = val
                    sf(cr, 13, False, _C(theme["dark"]))
        else:
            hint = slide.shapes.add_textbox(Inches(3.0), Inches(3.5), Inches(7.0), Inches(1.0))
            hr = hint.text_frame.paragraphs[0]; hr.alignment = PP_ALIGN.CENTER
            hrr = hr.add_run(); hrr.text = "[ 对比数据待补充 ]"; sf(hrr, 18, None, _C(theme["gray"]))
        _add_page_number(slide, total, idx, theme)

    # ---- summary ----
    def render_summary(s, idx):
        slide = ns()
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12.0), Inches(1.0))
        r = tb.text_frame.paragraphs[0].add_run(); r.text = s.get("title", "总结")
        sf(r, 30, True, _C(theme["primary"]))
        rect(slide, Inches(0.7), Inches(1.35), Inches(0.08), Inches(0.5), theme["accent"])
        rect(slide, Inches(0.85), Inches(1.37), Inches(3.5), Pt(3.5), theme["secondary"])

        body_tb = slide.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(11.5), Inches(3.8))
        btf = body_tb.text_frame; btf.word_wrap = True
        for i, pt in enumerate(s.get("points", [])):
            p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            p.space_before = Pt(8); p.space_after = Pt(6)
            cr = p.add_run(); cr.text = "\u2713  "; sf(cr, 16, False, _C(theme["success"]))
            tr = p.add_run(); tr.text = pt; sf(tr, 18, False, _C(theme["dark"]))

        conclusion = s.get("conclusion", "")
        if conclusion:
            rect(slide, Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.4), theme["light_gray"])
            cbt = slide.shapes.add_textbox(Inches(1.3), Inches(5.65), Inches(10.7), Inches(1.0))
            cbp = cbt.text_frame.paragraphs[0]; cbp.alignment = PP_ALIGN.CENTER
            cbr = cbp.add_run(); cbr.text = conclusion; sf(cbr, 22, True, _C(theme["primary"]))
        _add_page_number(slide, total, idx, theme)

    # ---- dispatch ----
    DISPATCH = {
        "cover": render_cover, "agenda": render_agenda, "section": render_section,
        "content": render_content, "two_column": render_two_column, "chart": render_chart,
        "table": render_table, "quote": render_quote, "timeline": render_timeline,
        "comparison": render_comparison, "summary": render_summary,
    }

    for idx, s in enumerate(slide_list, 1):
        fn = DISPATCH.get(s.get("type"), render_content)
        try:
            fn(s, idx)
        except Exception:
            # Fallback: ensure every slide renders something readable
            fs = ns()
            ftb = fs.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12.0), Inches(1.0))
            fr = ftb.text_frame.paragraphs[0].add_run()
            fr.text = s.get("title", "内容"); sf(fr, 28, True, _C(theme["primary"]))
            fbody = fs.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(11.5), Inches(5.0))
            fbf = fbody.text_frame; fbf.word_wrap = True
            fbr = fbf.paragraphs[0].add_run()
            items = s.get("items", []) or [s.get("text", "(内容加载异常)")]
            fbr.text = "\n".join("\u2022 " + it for it in items)
            sf(fbr, 18, False, _C(theme["dark"]))
            _add_page_number(fs, total, idx, theme)

    return prs


def _build_and_save(sd, out_path):
    prs = _render_pptx(sd)
    prs.save(out_path)
    return out_path


# ============================================================
#  Main entry point
# ============================================================
def _auto_fill_content_from_history(ctx, topic):
    """P2⑮ 入口兜底：当 LLM 调 generate_ppt 时 content 为空、且 ctx.messages 携带聊天历史时，
    自动从最近一条 assistant 消息抽取资料作为 content 填入。

    设计动机：工具是「无状态的执行器」——它不知道上下文里有什么数据。
    即使 system prompt 明确要求"必须传 content"，LLM 在反思/重试路径中
    仍可能只传 topic 不传 content，PPT 沦为通用空模板。
    兜底策略：当 ctx.history 含有"看起来像资料"的 assistant 文本（>=80 字符、
    非"工具执行成功"类响应）时，自动作为 content 注入。

    返回：(auto_content, note) — auto_content 为自动填入的 content，note 为
    写进返回文本的提示。auto_content 为 None 表示不注入。
    """
    if not ctx or not getattr(ctx, "messages", None):
        return None, ""
    _msgs = ctx.messages or []
    # 倒序找最近一条 assistant 文本（跳过 system/tool 角色）
    _last_assistant = None
    for m in reversed(_msgs):
        if isinstance(m, dict) and m.get("role") == "assistant":
            _content = (m.get("content") or "").strip()
            # 跳过"工具执行结果回灌"和"已生成 X"类的过短/过结构化文本
            # 额外过滤：噪声密集（连续 5+ 行 ≤2 字符的 LLM 反思/规划文本，注入会导致 PPT 破折号爆炸）
            if (len(_content) >= 80 and
                    "已生成：" not in _content and
                    "工具执行" not in _content and
                    not _is_noise_dense(_content)):
                _last_assistant = _content
                break
    if not _last_assistant:
        return None, ""
    # 截断保护（assistant 文本可能很长；上限 6000 字符够用）
    _auto = _last_assistant[:6000]
    _note = ("\n（⚠️ 你本次未传 content，已自动从聊天历史中提取最近一条助手消息作为 PPT 资料，"
             "下次请显式把要写入 PPT 的全部内容填入 content 参数，以获得更高质量的输出）")
    return _auto, _note


async def run(ctx, topic, content="", style="business_blue", slides=None):
    # P2 阶段 2：优先使用结构化 slides JSON 入参（LLM 直接传结构数组，绕开 Markdown 启发式解析）
    # 传 content 则走现有 _parse_content 启发式解析（向后兼容）
    if slides and isinstance(slides, list) and len(slides) > 0:
        _sd = _build_slides_from_json(topic, slides, style)
        if _sd:
            slides = _sd
            _use_json = True
        else:
            _use_json = False
    else:
        _use_json = False
    topic = (topic or "").strip()
    if not topic:
        return "PPT 生成失败：请提供 PPT 主题（topic 不能为空）。"

    style = (style or _DEFAULT_STYLE).strip()
    if style not in THEMES:
        style = _DEFAULT_STYLE

    # P2⑮ 入口兜底：content 为空时自动从 ctx.messages 抽取最近一条 assistant 文本作为 content。
    # 解决"LLM 反思重试只传 topic 不传 content"导致的凭空编造/通用空模板问题。
    _auto_note = ""
    if not (content or "").strip():
        _auto_content, _auto_note = _auto_fill_content_from_history(ctx, topic)
        if _auto_content:
            content = _auto_content

    # ---- 去重检查：同一 topic+当天 24h 内复用已有文件，避免闭环反思/反复调用重复生成 ----
    # P2⑮ 鲁棒化：用 _ppt_cache_key() 归一化 topic + 加当天日期过滤。
    # P2 用户隔离：cache_key 加 user_id 前缀，防止不同用户复用彼此的去重缓存（变相泄密）。
    # 例外：若「已缓存的是无内容默认模板」而「本次带真实 content」，则放行重新渲染
    # （用真实内容覆盖），这样 Agent 先调空内容、再补内容的常见流程仍能得到有内容的版本，
    # 且最多多生成 1 个文件。
    _uid = _safe_id(ctx.user.get("id") if getattr(ctx, "user", None) else None)
    cache_key = "%s::%s" % (_uid, _ppt_cache_key(topic))
    now = time.time()
    _has_content = bool((content or "").strip())
    cached = _recent_ppt_cache.get(cache_key)
    if cached and (now - cached[0]) < _CACHE_TTL:
        _cached_path, _cached_had_content = cached[1], cached[2]
        if os.path.exists(_cached_path):
            if _has_content and not _cached_had_content:
                pass  # 本次有内容而缓存的是空模板 → 放行重新渲染，覆盖默认模板
            else:
                return ("PPT 已生成（复用已有文件，风格：%s）：%s\n"
                        "（同一主题在 24h 内已生成过，直接返回已有文件以避免重复生成）" % (
                            THEMES.get(style, {}).get("name", style), _cached_path))
    # 清理过期缓存（原地删除，避免函数内重新赋值触发局部变量语义问题）
    _expired = [k for k, v in _recent_ppt_cache.items()
                if now - v[0] >= _CACHE_TTL * 2]
    for _k in _expired:
        del _recent_ppt_cache[_k]

    # slides 来源：函数顶部已处理 JSON 模式（_use_json=True）；此处仅 content 模式走启发式解析
    if not _use_json:
        slides = _parse_content(topic, content, style)
        if not slides:
            slides = _default_slides(topic, style)
    else:
        if not slides or not slides.get("slides"):
            slides = _default_slides(topic, style)

    # 稀烂 content 回退默认模板：有效内容页<2 时直接走默认模板（4 页+封面），
    # 避免「一行一页 117 页全是空破折号」式废稿。
    if _count_effective_content(slides) < 2:
        slides = _default_slides(topic, style)
    slide_list = slides.get("slides", [])

    # 总页数硬上限 20 页：超过截断并告警（建议拆分为多份或多套配色，或改为 Word 报告交付）
    _MAX_PAGES = 20
    _truncated = False
    if len(slide_list) > _MAX_PAGES:
        slides["slides"] = slide_list[:_MAX_PAGES]
        slide_list = slides["slides"]
        _truncated = True

    root = _user_dir(_uid, getattr(ctx, "session_id", None), "artifacts")
    try:
        os.makedirs(root, exist_ok=True)
    except Exception:
        root = tempfile.gettempdir()

    safe_name = re.sub(r"\W+", "_", topic[:30]) or "PPT"
    out = os.path.join(root, "%s_%s.pptx" % (safe_name, uuid.uuid4().hex[:8]))

    try:
        await asyncio.to_thread(_build_and_save, slides, out)
    except Exception as e:
        return "PPT 生成失败：%s: %s" % (type(e).__name__, e)

    theme_display = THEMES.get(style, {}).get("name", style)
    # 写入去重缓存（记录是否含真实内容，供后续"空模板→有内容"放行判断）
    _recent_ppt_cache[cache_key] = (now, out, _has_content)
    _trunc_note = ""
    if _truncated:
        _trunc_note = ("\n（⚠️ PPT 页数超过 %d 上限，已截断为 %d 页；如需更长内容，建议拆分为多份、"
                       "采用多套配色，或改为 Word 报告交付）" % (_MAX_PAGES, _MAX_PAGES))
    return ("PPT 已生成（%s风格，共 %d 页%s）：%s\n"
            "已生成：%s\n"
            "（已保存为 .pptx 文件，可在下方下载卡片中下载）%s%s" % (
                theme_display, len(slide_list),
                ("，已截断" if _truncated else ""),
                out, out, _auto_note, _trunc_note))
